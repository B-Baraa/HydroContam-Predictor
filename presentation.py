from pptx import Presentation
from pptx.util import Inches

# =========================
# CREATE PRESENTATION
# =========================
prs = Presentation()

# =========================
# FUNCTION
# =========================
def add_slide(title, content):
    layout = prs.slide_layouts[1]
    slide = prs.slides.add_slide(layout)
    slide.shapes.title.text = title
    slide.placeholders[1].text = content

# =========================
# TITLE SLIDE
# =========================
slide_layout = prs.slide_layouts[0]
slide = prs.slides.add_slide(slide_layout)
slide.shapes.title.text = "AI-Based Water Contamination Risk Prediction"
slide.placeholders[1].text = "Hybrid Scientific + Machine Learning Approach"

# =========================
# PROBLEM STATEMENT
# =========================
add_slide(
    "Problem Statement",
    """Groundwater contamination in Tunisia is increasing due to:
- Excessive use of chemical fertilizers
- Nitrate infiltration
- Microbial pollution

Goal:
Predict whether water is SAFE or HIGH RISK using data science."""
)

# =========================
# OBJECTIVES
# =========================
add_slide(
    "Project Objectives",
    """- Analyze water quality data
- Detect contamination risk
- Build a predictive ML model
- Create a scientific scoring system
- Develop an interactive dashboard"""
)

# =========================
# DATASETS
# =========================
add_slide(
    "Datasets Used",
    """- Groundwater Nitrate Dataset (Results.csv)
- Water Potability Dataset
- Water Quality Dataset

Each dataset contains chemical, physical, and biological indicators."""
)

# =========================
# DATA FUSION
# =========================
add_slide(
    "Data Fusion",
    """Multiple datasets were merged into a unified dataset.

Steps:
- Align features
- Select numeric attributes
- Concatenate datasets

Result:
One fused dataset for modeling"""
)

# =========================
# DATA CLEANING
# =========================
add_slide(
    "Data Preprocessing",
    """- Removed duplicates
- Handled missing values (mean imputation)
- Standardized column names
- Selected relevant features"""
)

# =========================
# DATA EXPLORATION
# =========================
add_slide(
    "Data Exploration",
    """Visualizations used:
- Histograms → distribution of features
- KDE plots → density analysis
- Pairplots → relationships between variables
- Violin plots → feature impact on risk

Goal:
Understand patterns and anomalies"""
)

# =========================
# IMPORTANT INSIGHT
# =========================
add_slide(
    "Key Insights",
    """- Nitrate strongly influences contamination
- Bacterial indicators (coliform) are critical
- Conductivity reflects dissolved pollutants
- Data contains many outliers → real-world variability"""
)

# =========================
# OUTLIER ANALYSIS
# =========================
add_slide(
    "Outlier Analysis",
    """Outliers detected using:
- Boxplots
- IQR method

Observation:
Outliers were NOT removed because they represent
real environmental extremes."""
)

# =========================
# FEATURE ENGINEERING
# =========================
add_slide(
    "Feature Engineering",
    """Created a contamination score using:

- Nitrate
- Conductivity
- WQI
- Coliform bacteria
- Turbidity
- Chloramines

Normalized and combined into a global score"""
)

# =========================
# SCIENTIFIC APPROACH
# =========================
add_slide(
    "Scientific Risk Model",
    """Instead of using only ML:

WHO standards were used:
- Nitrate > 50 mg/L → unsafe
- pH outside 6.5–8.5 → unsafe
- Coliform presence → unsafe

A rule-based scoring system improves reliability."""
)

# =========================
# TARGET VARIABLE
# =========================
add_slide(
    "Target Variable",
    """Risk label created using:

contamination_score > threshold → HIGH RISK
otherwise → SAFE

This ensures balanced classification."""
)

# =========================
# MODELING
# =========================
add_slide(
    "Machine Learning Models",
    """Models used:
- Random Forest
- Logistic Regression
- Decision Tree

Goal:
Compare performance and select best model."""
)

# =========================
# MODEL EVALUATION
# =========================
add_slide(
    "Model Evaluation",
    """Metrics used:
- Accuracy
- Confusion Matrix
- ROC Curve (AUC)
- Classification Report

Random Forest showed best performance."""
)

# =========================
# MODEL SELECTION
# =========================
add_slide(
    "Best Model",
    """Random Forest selected because:
- High accuracy
- Handles non-linear relationships
- Robust to noise and outliers"""
)

# =========================
# OVERFITTING CHECK
# =========================
add_slide(
    "Overfitting Analysis",
    """Compared:
- Train accuracy
- Test accuracy

Result:
Model generalizes well with minimal overfitting."""
)

# =========================
# FEATURE IMPORTANCE
# =========================
add_slide(
    "Feature Importance",
    """Most important variables:
- Nitrate
- Coliform bacteria
- Conductivity
- WQI

These directly impact water safety."""
)

# =========================
# SIMULATION
# =========================
add_slide(
    "Simulation",
    """Simulated contamination scenarios by:
- Increasing nitrate levels
- Observing model predictions

Result:
Risk increases with pollution levels."""
)

# =========================
# DASHBOARD
# =========================
add_slide(
    "Streamlit Dashboard",
    """Features:
- User input simulation
- Real-time prediction
- Scientific + ML analysis
- Visual feedback

Simulates IoT-based monitoring system."""
)

# =========================
# INNOVATION
# =========================
add_slide(
    "Innovation",
    """Hybrid System:
1. Scientific rule-based analysis
2. Machine learning prediction

This ensures:
- Explainability
- Accuracy
- Real-world applicability"""
)

# =========================
# LIMITATIONS
# =========================
add_slide(
    "Limitations",
    """- Limited real-time data
- Synthetic feature assumptions
- Dataset inconsistencies
- No live IoT integration yet"""
)

# =========================
# FUTURE WORK
# =========================
add_slide(
    "Future Work",
    """- Integrate IoT sensors
- Use real-time environmental data
- Add weather API
- Deploy on cloud
- Build mobile application"""
)

# =========================
# CONCLUSION
# =========================
add_slide(
    "Conclusion",
    """This project demonstrates how AI + environmental science
can help monitor water contamination.

The system is:
- Scalable
- Explainable
- Suitable for real-world deployment"""
)

# =========================
# SAVE FILE
# =========================
prs.save("water_contamination_full_presentation.pptx")

print("✅ Full professional presentation generated!")